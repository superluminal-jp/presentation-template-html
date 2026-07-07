"""HTML layout -> PowerPoint slide layout / placeholder mapping.

Maps the 16 `slides/NN-*.html` layouts (data-model.md "LayoutDefinition") to
PowerPoint custom slide layouts, and maps each extracted `SlideElement.role`
to a PPTX placeholder type (data-model.md §5 PlaceholderMapping).
"""

from copy import deepcopy

from pptx.opc.constants import CONTENT_TYPE as CT
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml.ns import qn
from pptx.parts.slide import SlideLayoutPart
from pptx.util import Emu

# --- Layout ID -> Japanese display name (FR-012) ---------------------------
# Order matches slides/NN-*.html and User Story 1's acceptance criteria.
LAYOUT_NAMES = {
    "01-title": "表紙",
    "02-toc": "目次",
    "03-section": "章扉",
    "04-body": "本文",
    "05-compare": "比較",
    "06-chart": "図表",
    "07-summary": "まとめ",
    "08-reference": "リファレンス",
    "09-big-number": "ビッグナンバー",
    "10-dashboard": "ダッシュボード",
    "11-timeline": "タイムライン",
    "12-matrix": "マトリクス",
    "13-process": "プロセス",
    "14-quote": "引用",
    "15-image-full": "全面画像",
    "16-closing": "クロージング",
}

LAYOUT_ORDER = list(LAYOUT_NAMES.keys())

# --- SlideElement.role -> OOXML placeholder `type` attribute ---------------
# (data-model.md §5 PlaceholderMapping). "other" has no entry: rendered as a
# plain (non-placeholder) text box.
PLACEHOLDER_PH_TYPE = {
    "title": "title",
    "subtitle": "subTitle",
    "body": "body",
    "caption": "body",
    "meta": "body",
    "image": "pic",
    "chart": "pic",
}

# Canvas the HTML is authored against (px), per docs/requirements.md 16:9 grid.
CANVAS_WIDTH_PX = 1280
CANVAS_HEIGHT_PX = 720

# Output slide size (FR-003).
SLIDE_WIDTH_EMU = Emu(int(13.333 * 914400))
SLIDE_HEIGHT_EMU = Emu(int(7.5 * 914400))


def bbox_ratio_to_emu(bbox_ratio: dict) -> dict:
    """Convert a 0-1 ratio bbox to absolute EMU offsets on the output canvas."""
    return {
        "left": Emu(int(bbox_ratio["x"] * SLIDE_WIDTH_EMU)),
        "top": Emu(int(bbox_ratio["y"] * SLIDE_HEIGHT_EMU)),
        "width": Emu(int(bbox_ratio["w"] * SLIDE_WIDTH_EMU)),
        "height": Emu(int(bbox_ratio["h"] * SLIDE_HEIGHT_EMU)),
    }


def px_bbox_to_ratio(bbox_px: dict) -> dict:
    """Convert a px bbox (against CANVAS_WIDTH_PX/HEIGHT_PX) to a 0-1 ratio bbox."""
    return {
        "x": bbox_px["x"] / CANVAS_WIDTH_PX,
        "y": bbox_px["y"] / CANVAS_HEIGHT_PX,
        "w": bbox_px["w"] / CANVAS_WIDTH_PX,
        "h": bbox_px["h"] / CANVAS_HEIGHT_PX,
    }


def _blank_base_layout(prs):
    """Return the stock 'Blank' slide layout used as a cloning base."""
    master = prs.slide_masters[0]
    return next(layout for layout in master.slide_layouts if layout.name == "Blank")


def _next_shape_id(spTree) -> int:
    ids = [int(cNvPr.get("id")) for cNvPr in spTree.findall(f".//{qn('p:cNvPr')}")]
    return (max(ids) if ids else 1) + 1


def clone_layout(prs, layout_name: str):
    """Create a new, empty custom slide layout named *layout_name*.

    Clones the stock 'Blank' layout (fewest pre-existing placeholders) and
    strips its date/footer/slide-number boilerplate shapes, leaving a bare
    spTree ready for our own placeholders to be added.
    """
    master = prs.slide_masters[0]
    base_layout = _blank_base_layout(prs)
    new_element = deepcopy(base_layout._element)

    cSld = new_element.find(qn("p:cSld"))
    cSld.set("name", layout_name)

    spTree = new_element.find(f"{qn('p:cSld')}/{qn('p:spTree')}")
    for sp in spTree.findall(qn("p:sp")):
        spTree.remove(sp)

    package = master.part.package
    partname = package.next_partname("/ppt/slideLayouts/slideLayout%d.xml")
    new_part = SlideLayoutPart(partname, CT.PML_SLIDE_LAYOUT, package, new_element)
    # Every slideLayout part must carry a relationship back to its slideMaster
    # (ECMA-376); python-pptx's SlideLayoutPart.slide_master looks this up via
    # part_related_by(RT.SLIDE_MASTER). Without it PowerPoint treats the
    # package as corrupt and either refuses to open it or silently repairs
    # (drops) the offending layout.
    new_part.relate_to(master.part, RT.SLIDE_MASTER)

    rId = master.part.relate_to(new_part, RT.SLIDE_LAYOUT)
    sldLayoutIdLst = master.part._element.find(qn("p:sldLayoutIdLst"))
    existing_ids = [int(e.get("id")) for e in sldLayoutIdLst.findall(qn("p:sldLayoutId"))]
    new_id = max(existing_ids) + 1
    new_entry = sldLayoutIdLst._add_sldLayoutId()
    new_entry.set("id", str(new_id))
    new_entry.rId = rId

    return new_part.slide_layout


_PLACEHOLDER_SP_TEMPLATE = """
<p:sp {nsdecls}>
  <p:nvSpPr>
    <p:cNvPr id="{shape_id}" name="{name}"{descr_attr}/>
    <p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>
    <p:nvPr><p:ph type="{ph_type}"{idx_attr}/></p:nvPr>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="{x}" y="{y}"/><a:ext cx="{cx}" cy="{cy}"/></a:xfrm>
  </p:spPr>
  <p:txBody>
    <a:bodyPr/>
    <a:lstStyle/>
    <a:p><a:endParaRPr lang="ja-JP"/></a:p>
  </p:txBody>
</p:sp>
"""


def add_placeholder_shape(layout, element: dict, idx: int) -> None:
    """Append an empty (prompt-text-only) placeholder shape for *element*.

    Text is deliberately left empty so PowerPoint shows its own built-in
    prompt text instead of any baked-in HTML copy (FR-004, Clarification #1).
    Alt text (FR-014) is set on the shape's cNvPr when `alt_text_hint` is
    present. Reading order (FR-015) follows insertion order in the spTree,
    which callers control via the `elements` iteration order (sorted by
    `reading_order_index`).
    """
    from pptx.oxml.ns import nsdecls

    role = element["role"]
    ph_type = PLACEHOLDER_PH_TYPE.get(role, "body")
    emu = bbox_ratio_to_emu(element["bbox_ratio"])
    spTree = layout._element.find(f"{qn('p:cSld')}/{qn('p:spTree')}")
    shape_id = _next_shape_id(spTree)

    needs_idx = ph_type != "title"
    idx_attr = f' idx="{idx}"' if needs_idx else ""

    alt_text = element.get("alt_text_hint")
    descr_attr = f' descr="{alt_text}"' if alt_text else ""

    xml = _PLACEHOLDER_SP_TEMPLATE.format(
        nsdecls=nsdecls("a", "p", "r"),
        shape_id=shape_id,
        name=f"{role.capitalize()} Placeholder {shape_id}",
        descr_attr=descr_attr,
        ph_type=ph_type,
        idx_attr=idx_attr,
        x=int(emu["left"]),
        y=int(emu["top"]),
        cx=int(emu["width"]),
        cy=int(emu["height"]),
    )
    from pptx.oxml import parse_xml

    sp = parse_xml(xml)
    spTree.append(sp)


def add_picture_shape(layout, element: dict, image_path: str):
    """Insert an actual raster Picture shape (fallback path, FR-008/FR-014).

    `LayoutShapes` (unlike `SlideShapes`) has no high-level `add_picture()` in
    python-pptx, so this replicates it at the part/oxml level: register the
    image part against the layout's part, then append a `p:pic` element to
    the layout's spTree directly.
    """
    emu = bbox_ratio_to_emu(element["bbox_ratio"])
    image_part, rId = layout.part.get_or_add_image_part(image_path)
    spTree = layout._element.find(f"{qn('p:cSld')}/{qn('p:spTree')}")
    shape_id = _next_shape_id(spTree)
    alt_text = element.get("alt_text_hint") or image_part.desc

    pic = spTree.add_pic(
        shape_id,
        f"Picture {shape_id}",
        alt_text,
        rId,
        int(emu["left"]),
        int(emu["top"]),
        int(emu["width"]),
        int(emu["height"]),
    )
    return pic
